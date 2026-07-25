#!/usr/bin/env python3
"""Remplace l'emplacement LOGO_CLIENT d'un PPTX par l'image d'un client.

    python3 logo_client_pptx.py kit.pptx acme.png -o kit-acme.pptx

L'emplacement est repéré par le NOM de la forme (LOGO_CLIENT), pas par sa
position : la mise en page peut évoluer sans casser l'automatisation.
Le logo est inscrit dans le cadre en respectant son rapport d'aspect, puis
centré ; le cadre pointillé et son libellé sont supprimés.
"""
import argparse
import pathlib
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

NOM = "LOGO_CLIENT"


def retirer(forme):
    forme._element.getparent().remove(forme._element)


def injecter(pptx_src, logo, pptx_out):
    pres = Presentation(pptx_src)
    lw, lh = Image.open(logo).size
    ratio = lw / lh
    poses = 0

    for i, slide in enumerate(pres.slides, 1):
        cibles = [s for s in slide.shapes if s.name == NOM]
        for cadre in cibles:
            cx, cy, cw, ch = cadre.left, cadre.top, cadre.width, cadre.height

            # Plaque blanche sous le logo : sur une couverture sombre, le logo
            # client (souvent sombre ou multicolore) ne ressort pas. On pose un
            # rectangle blanc arrondi à la taille du cadre ; le logo vient dessus
            # avec une marge intérieure pour respirer.
            plaque = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
            plaque.adjustments[0] = 0.10
            plaque.fill.solid()
            plaque.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            plaque.line.fill.background()
            plaque.shadow.inherit = False

            # inscription dans le cadre (avec marge), rapport d'aspect conservé
            pad_w, pad_h = int(cw * 0.10), int(ch * 0.16)
            iw, ih = cw - 2 * pad_w, ch - 2 * pad_h
            if iw / ih > ratio:
                h = ih
                w = int(ih * ratio)
            else:
                w = iw
                h = int(iw / ratio)
            slide.shapes.add_picture(str(logo), Emu(cx + (cw - w) // 2), Emu(cy + (ch - h) // 2),
                                     width=Emu(w), height=Emu(h))
            retirer(cadre)
            poses += 1
            print(f"  slide {i} : logo posé en {w/914400:.2f} × {h/914400:.2f} po")
        for txt in [s for s in slide.shapes if s.name == NOM + "_TEXTE"]:
            retirer(txt)

    if not poses:
        sys.exit(f"aucune forme nommée {NOM} dans {pptx_src}")
    pres.save(pptx_out)
    _reparer_content_types(pptx_out)
    print(f"→ {pptx_out}")


MIMES = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
         "gif": "image/gif", "svg": "image/svg+xml", "emf": "image/x-emf", "wmf": "image/x-wmf"}


def _reparer_content_types(chemin):
    """python-pptx réécrit [Content_Types].xml et peut omettre des extensions
    d'images présentes dans ppt/media. PowerPoint refuse alors le fichier."""
    import re
    import shutil
    import tempfile
    import zipfile

    with zipfile.ZipFile(chemin) as z:
        noms = z.namelist()
        ct = z.read("[Content_Types].xml").decode()
        contenus = {n: z.read(n) for n in noms}

    exts = {n.rsplit(".", 1)[-1].lower() for n in noms if n.startswith("ppt/media/") and "." in n}
    manquantes = [e for e in sorted(exts)
                  if e in MIMES and f'Extension="{e}"' not in ct]
    if not manquantes:
        return

    ajout = "".join(f'<Default Extension="{e}" ContentType="{MIMES[e]}"/>' for e in manquantes)
    ct = re.sub(r"(<Types[^>]*>)", r"\1" + ajout, ct, count=1)
    contenus["[Content_Types].xml"] = ct.encode()

    tmp = tempfile.mktemp(suffix=".pptx")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as z:
        for n in noms:
            z.writestr(n, contenus[n])
    shutil.move(tmp, chemin)
    print(f"  extensions déclarées : {', '.join(manquantes)}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("logo")
    ap.add_argument("-o", "--out", required=True)
    a = ap.parse_args()
    for f in (a.pptx, a.logo):
        if not pathlib.Path(f).exists():
            sys.exit(f"introuvable : {f}")
    injecter(a.pptx, a.logo, a.out)
